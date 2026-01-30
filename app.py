import streamlit as st
from openai import OpenAI
from pptx import Presentation
from io import BytesIO

# 1. Setup OpenAI Client
# This reads the key from Streamlit's Secret Manager
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

def generate_slide_content(topic, num_slides):
    prompt = (
        f"Create an outline for a {num_slides}-slide presentation on '{topic}'.\n"
        "Format each slide exactly like this:\n"
        "TITLE: [Slide Title]\n"
        "CONTENT: [Bullet point 1], [Bullet point 2], [Bullet point 3]"
    )
    
    response = client.chat.completions.create(
        model="gpt-4o",  # You can also use "gpt-3.5-turbo"
        messages=[
            {"role": "system", "content": "You are a professional presentation assistant."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content

def create_pptx(topic, content):
    prs = Presentation()
    
    # Add a Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "AI Generated with OpenAI & Python"

    # Split the AI response into individual slides
    slides_raw = content.split("TITLE:")
    for slide_data in slides_raw[1:]:  # Skip first empty split
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Split Title and Content
        parts = slide_data.split("CONTENT:")
        slide.shapes.title.text = parts[0].strip()
        
        if len(parts) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            bullets = parts[1].strip().split(",")
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet.strip()

    binary_output = BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()

# 2. Streamlit Interface
st.set_page_config(page_title="OpenAI Slide Maker", page_icon="📊")
st.title("📊 OpenAI Slide Generator")

topic = st.text_input("Presentation Topic:", placeholder="e.g. The Future of Renewable Energy")
num_slides = st.slider("Number of Slides", 3, 10, 5)

if st.button("Generate PowerPoint"):
    with st.spinner("OpenAI is drafting your slides..."):
        try:
            text_content = generate_slide_content(topic, num_slides)
            pptx_file = create_pptx(topic, text_content)
            
            st.success("Presentation Ready!")
            st.download_button(
                label="📥 Download .pptx",
                data=pptx_file,
                file_name=f"{topic.replace(' ', '_')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"Error: {e}")
